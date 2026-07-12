# -*- coding: utf-8 -*-
"""
build_ppt.py — 答辩 PPT 生成器骨架(python-pptx,16:9)
纪律与 build_report.py 相同:数字全部复用其 CSV 装载结果(import 即共享,单一来源),
叙事零硬编码性能数字;图直接嵌 sim/figs。演讲稿口径以本 PPT 页序为准(R2 培养体系配套)。
运行:python build_ppt.py → out/答辩PPT_draft.pptx
"""
import os
import sys
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

import build_report as br                      # 复用全部数字装载(不触发其 main)
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs")
FONT = "Microsoft YaHei"
NAVY = RGBColor(0x17, 0x36, 0x5D)
GRAY = RGBColor(0x50, 0x50, 0x50)

prs = Presentation()
prs.slide_width, prs.slide_height = Cm(33.87), Cm(19.05)      # 16:9
BLANK = prs.slide_layouts[6]
prs.core_properties.title = "智储优控答辩PPT（生成器草稿）"
prs.core_properties.subject = "sim口径与PLC证据边界分离；canonical v3.4"
prs.core_properties.keywords = "generated=" + datetime.now().isoformat(timespec="seconds")

_l4 = [r for r in br.load_rows("l4_task_cycles_data.csv") if r.get("cycles_to_done")]
try:
    _a = [r for r in _l4 if r["group"] == "A"][0]["cycles_to_done"]
    _c = [r for r in _l4 if r["group"] == "C"][0]["cycles_to_done"]
    l4_bullet = ("历史AB PC仿真:nBatch=1→%s周期;nBatch≥5约%s周期;×10ms仅为"
                 "标称任务周期换算，实体AC500最坏时延/watchdog待测" % (_a, _c))
except Exception:
    l4_bullet = "〔L4 数据读取失败,查 sim/out/l4_task_cycles_data.csv〕"


def _set(run, size, bold=False, color=None):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or RGBColor(0x20, 0x20, 0x20)


def slide(title, bullets=None, fig=None, note=""):
    """通用页:左题+要点列表,可选右侧/下方图;note=演讲提词(进备注栏)。"""
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Cm(1.2), Cm(0.7), Cm(31), Cm(1.8))
    _set(tb.text_frame.paragraphs[0].add_run(), 28, True, NAVY)
    tb.text_frame.paragraphs[0].runs[0].text = title
    if bullets:
        bx = s.shapes.add_textbox(Cm(1.4), Cm(3.0),
                                  Cm(14.5 if fig else 30.5), Cm(14.5))
        tf = bx.text_frame
        tf.word_wrap = True
        for i, (txt, lv) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.level = lv
            r = p.add_run()
            r.text = ("• " if lv == 0 else "– ") + txt
            _set(r, 16 if lv == 0 else 13.5,
                 bold=(lv == 0), color=None if lv == 0 else GRAY)
    if fig:
        path = os.path.join(FIGS, fig)
        if os.path.exists(path):
            s.shapes.add_picture(path, Cm(16.4 if bullets else 4.0), Cm(3.2),
                                 width=Cm(16.5 if bullets else 26.0))
    if note:
        s.notes_slide.notes_text_frame.text = note
    return s


# ---------------- 页序(≈18 页;〔〕为占位待人工补) ----------------
s = prs.slides.add_slide(BLANK)                                  # 1 封面
tb = s.shapes.add_textbox(Cm(2), Cm(6.2), Cm(30), Cm(6))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
_set(p.add_run(), 36, True, NAVY)
p.runs[0].text = "智储优控:立体仓储仓位优化\n基于 AI 算法与智能控制"
p2 = tb.text_frame.add_paragraph()
_set(p2.add_run(), 18, False, GRAY)
p2.runs[0].text = "2026 ABB 杯 · 赛道一 · 队伍信息提交前填写"

slide("为什么这道题不简单",
      [("400 仓位、承重逐层递减、预占用、五属性货物——多目标互相冲突", 0),
       ("存取效率 vs 重货稳定 vs 提升能耗 vs 近位稀缺:没有免费午餐", 1),
       ("分层交付:sim数据链 / ST核心逻辑 / 二进制制品分别验收", 0),
       ("主办方定调:打造“会思考”的立体仓储系统,绿色工程视野(王余,开赛致辞)", 1)],
      note="开场 40 秒:题面复述+难点定位,引用致辞建立同频。")

slide("双层优化架构:在线评分 + 批量 AWRA-LS",
      [("在线层:ST单call扫描400仓位；AB PC仿真单周期返回，实体目标机最坏时延待测", 0),
       ("四项加权:效率 + 稳定 + 能耗 + 位置价值预留 δ", 1),
       (f"批量层:Rank→Assign→LocalSearch,较在线再降 {br.GAP_ONLINE}%(信息完备的价值)", 0),
       (f"场景权重表fail-first:{br.ADW_INFEASIBLE}格INFEASIBLE、{br.ADW_UNSTABLE}格2-seed不稳；PLC未接入", 1)],
      note="架构页 60 秒:强调双层分别对应决赛交互与优化上界。")

slide("δ 预留项:一个有证据链的设计",
      [("δ=0仅证明占用仓位集合与near一致，不证明gid映射/严格算法等价", 0),
       ("低填充生命周期drain重算:δ=0.15显著更慢", 0),
       ("240件仅见失败缓解但仍大量失败——不称可行性保险", 1)],
      fig="fig14_参数敏感性.png",
      note="评委最容易问'权重怎么定的'——这页就是答案。")

slide("物理与作业真实度",
      [("梯形加减速运动模型(匀速口径保留为文献对照)", 0),
       (f"存取混合双命令周期:较单命令节省 {br.DUAL_SAVE}%", 0),
       (f"吞吐{br.THR_AW}单操作/时(约{br.THR_CYCLES}双周期/时)=行程-only连续满载sim容量", 1),
       ("不含装卸/排队/故障/扫描，不是PLC实测", 1)],
      fig="fig8_混合作业.png",
      note="强调'贴近实际工程应用'是题面原话。")

slide("主结果(sim H设计口径,5-seed回归锚)",
      [(f"AWRA-LS {br.H_AW}±{br.H_AW_STD}s(sample SD;95%CI半宽{br.H_AW_CI}),较基线降{br.DROP_VS_SEQ}%", 0),
       (f"能耗代理降{br.DROP_E}%;重货平均层{br.HEAVY_TIER};失败0/承重容积复查0", 0),
       ("非U3 30-seed泛化/非独立final；不是PLC H闭环", 1),
       (f"对照口径(匀速+固定权重,可比文献):{br.LEG_AW}s / 降 {br.LEG_DROP}%", 1)],
      fig="fig1_策略对比.png",
      note="全场最重的一页,放慢;口径两句话讲清,体现严谨。")

slide("离下界与已知可达点有多远",
      [(f"静态匀速sim夹逼:{br.LB_S}s≤可达点{br.ACH_S}s≤AWRA {br.AW_BOUND_S}s", 0),
       (f"距已证可达点{br.GAP_ACH}%（不是理论天花板；与H 8.40s不是同一物理口径）", 0),
       ("随机策略解析值 vs 仿真:1σ 内互证,仿真器本身被理论校验", 1)],
      fig="fig9_理论对比.png",
      note="三级证据链是名校队语言;一句话金句放这页。")

slide("一套系统,九把尺子:口径矩阵(方法论)",
      [("每个建模选择都被识别为“口径”:参数化+保留切换开关+主/对照并列呈现", 0),
       ("运动模型(加减速/匀速)/ 权重(自适应/固定)/ 行程(切比雪夫/曼哈顿)…共 9 维", 1),
       ("证据边界:120件H与AWRA-vs-seq行程方向保持；高密度/生命周期有负结果", 0),
       ("类比:汽车油耗的 WLTP vs EPA、财报的双准则——工程数字必须声明测量框架", 1),
       ("完整矩阵见报告附录 C;主口径数字由回归测试锁定,动口径=红灯", 1)],
      note="创新10%+报告30%的方法论页;评委问'为什么两套数'时,这页是总答案。")

slide("拟真度阶梯:三档sim同管线验证(L7)",
      [("档1 数据模型→档2 sim H设计口径→档3 综合sim场景", 0),
       ("档3 = 共同泊松到达流+装卸+频次噪声+故障；参数为显式假设", 1),
       (f"降幅阶梯:{br.TIER_DROP['1']}% → {br.TIER_DROP['2']}% → {br.TIER_DROP['3']}%"
        "——AWRA相对seq行程方向保持；失败/承重容积复查0", 0),
       (f"档3 排队视角:响应 P95 {br.RESP_SEQ95}s → {br.RESP_AW95}s"
        f"；near P95={br.RESP_NEAR95}s < AWRA={br.RESP_AW95}s，不能称AWRA全面最优", 1),
       ("频次画像带噪(σ=0.3)AWRA 布局质量不塌:多因子 Rank+局部搜索兜底", 1)],
      fig="fig16_拟真度阶梯.png",
      note="共同到达流由seq reference在策略循环外生成一次；逐指标解读。")

slide("鲁棒性与运行期自愈",
      [("F12特定3-seed压力场景:AWRA系退化较缓、该紧库存样本失败0件", 0),
       ("sim周期重排机制:消除布局熵漂移(收益>1.2×成本才动,经济学准入)", 0),
       ("负结果如实呈现:重排的价值是自愈不是提效(F11)", 1)],
      fig="fig13_鲁棒压力.png",
      note="'负结果也如实写'是工程可信度的加分句。")

slide("AC500 ST实现:有限AB PC仿真证据",
      [("三层封装:FC纯函数/FB状态机/PRG主状态机", 0),
       ("扫描周期分片架构；400件最坏时延/watchdog仍待目标机验证", 0),
       (f"已验收基线:核心ST+静态匀速黄金向量AB PC仿真{br.TEST_N}/{br.TEST_FAIL}", 0),
       (f"当前源码{br.SOURCE_TEST_N}项:{br.SOURCE_REAL_CASES}真实断言、{br.PLACEHOLDER_CASES}占位；未AB复测", 1),
       ("CB/TOB分配与lex路由已写入源码；源码存在不升级59/0证据", 1),
       ("三时钟源(LTIME/Monitor/SysTimeGetNs)周期内耗时在 PC 仿真均恒 0——如实记录,口径转任务级", 1),
       (l4_bullet, 1)],
      note="明确区分AB PC仿真、标称周期换算与实体目标机实测。")

slide("有限一致性护栏:Python 与 ST",
      [("20件T25静态匀速黄金向量逐位一致", 0),
       ("near平局键已在源码统一并增加T60；尚未AB复测", 0),
       ("仍知差异:REAL精度、加减速评分分母", 0),
       ("59/0不等于报告H KPI在AC500复现", 0),
       ("一致性探针不模拟全部PLC状态/float32路径", 1)],
      note="只讲已锁静态向量和已知差异，不把有限一致性升级成数字孪生闭环。")

slide("绿色工程:能耗的年化账",
      [(f"sim净机械功情景换算:年省约{br.KWH_SAVE}kWh(降{br.KWH_SAVE_PCT}%)", 0),
       (f"{br.DAILY_OPS}单操作/日={br.DAILY_DUAL}双周期；CO2因子{br.CO2_FACTOR}(2023全国平均)", 1),
       ("未计待机/辅助/驱动器/动态损耗；非电表或PLC实测", 1),
       ("呼应“AI 算法正在重新定义能效的边界”(王余)", 0)],
      fig="fig10_绿色换算.png",
      note="封面级卖点;强调假设显式=可质疑可替换。")

slide("决赛可视化演示(录屏/现场)",
      [("三页制706元素结构已生成；输入动作GUI洗+在线长跑仍待验证", 0),
       ("演示脚本v4:12个主镜头+4个补充镜头；逐镜头只讲画面实际可证状态", 0),
       ("δ调参仅作静态敏感性演示；δ=0不等价于near，现场联动效果须用户在场验收", 0),
       ("当前页不嵌未完成录屏；正式截图须来自用户在场的在线长跑验收", 1)],
      note="决赛 20% 的正面战场;留出现场操作时间。")

slide("创新点清单",
      [("δ自我修正证据链(存在性→标定→lifecycle反证与边界收窄)", 0),
       ("场景权重fail-first可行性/稳定性表（负结果透明；PLC未闭环）", 0),
       ("扫描周期分片状态机(优化算法的 PLC 原生工程化)", 0),
       ("Python↔ST有限一致性护栏；sim周期重排机制", 0)],
      note="每条都有实验编号与复现命令,防'创新是嘴上说的'质疑。")

slide("工程取舍(考虑过且不用,都有数据)",
      [("A* 路径规划:双轴同动+无障碍,切比雪夫直达即物理最优", 0),
       (f"强化学习:**自己实跑了轻量版**——自学 1200 回合到 {br.RL_FINAL}s,仍落后手工标定 "
        f"{br.RL_GAP}%,且未形成δ结构(余弦 {br.RL_COS})；结合生命周期负结果，不把它当δ价值证据", 0),
       ("深度版引 2025 文献:需长期数据、收益个位数、黑箱——不符合实时可审计要求", 1),
       ("查表化执行:在线评分本身微秒级,查表反丢解释链", 0)],
      fig="fig15_RL对照.png",
      note="'知道为什么不用'且'亲手试过'——这页现在有自己的学习曲线。")

slide("结论",
      [(f"sim H回归锚:取货降{br.DROP_VS_SEQ}%、能耗代理降{br.DROP_E}%（5-seed sample SD）", 0),
       (f"PLC证据边界:{br.TEST_N}/{br.TEST_FAIL}静态匀速黄金向量；不等于H KPI闭环", 0),
       ("分阶段复现:core_smoke/report_rebuild/artifact_verify；长专项与AB另验", 0)],
      note="收尾引用张楠'从学子到工程师的思维跃迁'致谢即可。")

slide("Q&A", [("口径优先级:canonical → 当前CSV → 历史证据；sim与PLC实测分开回答", 0)])

path = os.path.join(HERE, "out", "答辩PPT_draft.pptx")
prs.save(path)
print(f"已生成:{path}({len(prs.slides.__iter__.__self__._sldIdLst)} 页)")
print("生成完成；队伍信息与正式在线截图须在提交前由用户确认，不以占位符伪装完成。")
